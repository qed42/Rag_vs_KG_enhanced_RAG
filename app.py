import streamlit as st
import os
from dotenv import load_dotenv
from PyPDF2 import PdfReader
import docx2txt
from pptx import Presentation
import pandas as pd
from pathlib import Path
import tempfile
import json
import time
from datetime import datetime
import networkx as nx
from pyvis.network import Network
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAI, OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.prompts import PromptTemplate
from langchain.chains import LLMChain, RetrievalQA
from neo4j import GraphDatabase
from difflib import SequenceMatcher

def display_chat_history():
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []

    for i, (question, trad_response, kg_response) in enumerate(st.session_state.chat_history):
        st.markdown("---")
        st.markdown(f"### Question {i+1}")
        st.info(f"Q: {question}")
        
        col1, col2 = st.columns(2)
        with col1:
            st.markdown("#### Traditional RAG Response")
            st.success(trad_response[0])
            st.metric("Response Time", f"{trad_response[1]:.2f}s")
        
        with col2:
            st.markdown("#### KG-Enhanced RAG Response")
            st.success(kg_response[0])
            st.metric("Response Time", f"{kg_response[1]:.2f}s")

# Load environment variables
load_dotenv()

# Configure Streamlit page
st.set_page_config(
    page_title="Enhanced RAG System",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS
st.markdown("""
    <style>
    .stApp {
        background-color: #1E1E1E;
        color: white;
    }
    .main-status {
        background-color: #2D2D2D;
        border-radius: 10px;
        padding: 15px;
        margin: 10px 0;
        border: 2px solid #4CAF50;
    }
    .status-card {
        background-color: #3D3D3D;
        color: white !important;
        border-radius: 8px;
        padding: 15px;
        margin: 10px 0;
        border: 1px solid #4CAF50;
    }
    .graph-card {
        background-color: #2D2D2D;
        border-radius: 10px;
        padding: 20px;
        margin: 10px 0;
        border: 1px solid #4CAF50;
    }
    .file-list {
        background-color: #2D2D2D;
        padding: 10px;
        border-radius: 5px;
        margin-top: 5px;
    }
    .stButton>button {
        background-color: #4CAF50;
        color: white;
        border-radius: 5px;
        border: none;
        padding: 10px 24px;
        width: 100%;
    }
    .stTextInput>div>div>input {
        background-color: #3D3D3D;
        color: white;
    }
    </style>
""", unsafe_allow_html=True)

class DocumentProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )

    def read_pdf(self, file):
        pdf_reader = PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text

    def read_docx(self, file):
        return docx2txt.process(file)

    def read_txt(self, file):
        return file.read().decode('utf-8')

    def read_pptx(self, file):
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def process_file(self, file):
        file_extension = Path(file.name).suffix.lower()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
            tmp_file.write(file.getvalue())
            tmp_file_path = tmp_file.name

        try:
            if file_extension == '.pdf':
                text = self.read_pdf(tmp_file_path)
            elif file_extension == '.docx':
                text = self.read_docx(tmp_file_path)
            elif file_extension == '.txt':
                text = self.read_txt(file)
            elif file_extension == '.pptx':
                text = self.read_pptx(tmp_file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_extension}")

            chunks = self.text_splitter.split_text(text)
            return chunks

        finally:
            os.unlink(tmp_file_path)

class KnowledgeGraphManager:
    def __init__(self, uri, user, password):
        self.driver = GraphDatabase.driver(uri, auth=(user, password))
        
    def close(self):
        self.driver.close()
        
    def clear_database(self):
        with self.driver.session() as session:
            session.run("MATCH (n) DETACH DELETE n")

    def create_graph_visualization(self):
        G = nx.Graph()
        
        with self.driver.session() as session:
            nodes = session.run("""
                MATCH (n)
                RETURN DISTINCT labels(n) as label, n.name as name, id(n) as id
            """)
            
            relationships = session.run("""
                MATCH (a)-[r]->(b)
                RETURN id(a) as source_id, id(b) as target_id,
                       a.name as source_name, b.name as target_name,
                       type(r) as relationship_type
            """)
            
            for node in nodes:
                G.add_node(node['id'], 
                          label=node['name'], 
                          title=f"{node['label'][0]}: {node['name']}")
            
            for rel in relationships:
                G.add_edge(rel['source_id'], 
                          rel['target_id'], 
                          title=rel['relationship_type'])
        
        net = Network(height="600px", width="100%", bgcolor="#222222", font_color="white")
        
        net.set_options("""
        {
          "physics": {
            "forceAtlas2Based": {
              "gravitationalConstant": -50,
              "springLength": 100,
              "springConstant": 0.08
            },
            "maxVelocity": 50,
            "solver": "forceAtlas2Based",
            "timestep": 0.35
          },
          "nodes": {
            "color": {
              "background": "#4CAF50",
              "border": "#2E7D32"
            }
          },
          "edges": {
            "color": {
              "color": "#FFFFFF",
              "highlight": "#4CAF50"
            }
          }
        }
        """)
        
        net.from_nx(G)
        net.save_graph("temp_graph.html")
        
        with open("temp_graph.html", "r", encoding="utf-8") as f:
            html_string = f.read()
            
        os.remove("temp_graph.html")
        return html_string

    def get_graph_statistics(self):
        with self.driver.session() as session:
            node_counts = session.run("""
                MATCH (n)
                WITH labels(n) as labels
                UNWIND labels as label
                RETURN label, count(*) as count
                ORDER BY count DESC
            """)
            
            rel_counts = session.run("""
                MATCH ()-[r]->()
                WITH type(r) as type
                RETURN type, count(*) as count
                ORDER BY count DESC
            """)
            
            totals = session.run("""
                MATCH (n)
                OPTIONAL MATCH (n)-[r]->()
                RETURN count(DISTINCT n) as nodes, count(DISTINCT r) as relationships
            """)
            
            return {
                'node_counts': [(record['label'], record['count']) for record in node_counts],
                'relationship_counts': [(record['type'], record['count']) for record in rel_counts],
                'totals': next(totals)
            }
    def create_graph_visualization(self):
        """Create an interactive graph visualization"""
        try:
            with self.driver.session() as session:
                # Get nodes
                nodes = session.run("""
                    MATCH (n)
                    RETURN id(n) as id, n.name as label, labels(n) as type
                """)
                
                # Get relationships
                edges = session.run("""
                    MATCH (a)-[r]->(b)
                    RETURN id(a) as source, id(b) as target, 
                           type(r) as label, r.type as relationship_type
                """)

                # Create network
                net = Network(
                    height="600px", 
                    width="100%", 
                    bgcolor="#222222", 
                    font_color="white",
                    directed=True
                )

                # Add nodes
                for node in nodes:
                    net.add_node(
                        node['id'],
                        label=node['label'],
                        title=f"Type: {node['type'][0]}",
                        color="#4CAF50"
                    )

                # Add edges
                for edge in edges:
                    net.add_edge(
                        edge['source'],
                        edge['target'],
                        label=edge['relationship_type'] or edge['label'],
                        color="#FFFFFF"
                    )

                # Set physics layout
                net.set_options("""
                    {
                        "physics": {
                            "forceAtlas2Based": {
                                "gravitationalConstant": -50,
                                "springLength": 100,
                                "springConstant": 0.08
                            },
                            "maxVelocity": 50,
                            "solver": "forceAtlas2Based",
                            "timestep": 0.35
                        },
                        "nodes": {
                            "font": {
                                "size": 12,
                                "color": "white"
                            }
                        },
                        "edges": {
                            "font": {
                                "size": 10,
                                "color": "white"
                            },
                            "arrows": {
                                "to": {
                                    "enabled": true,
                                    "scaleFactor": 0.5
                                }
                            }
                        }
                    }
                """)

                # Save and return HTML
                net.save_graph("temp_graph.html")
                with open("temp_graph.html", "r", encoding="utf-8") as f:
                    html_string = f.read()
                os.remove("temp_graph.html")
                return html_string

        except Exception as e:
            st.error(f"Error creating graph visualization: {str(e)}")
            return None

    def get_graph_summary(self):
        """Get a summary of the knowledge graph"""
        try:
            with self.driver.session() as session:
                summary = session.run("""
                    MATCH (n)
                    WITH labels(n) as label, count(n) as count
                    RETURN label, count
                    ORDER BY count DESC
                """)
                
                relationships = session.run("""
                    MATCH ()-[r]->()
                    WITH type(r) as type, count(r) as count
                    RETURN type, count
                    ORDER BY count DESC
                """)
                
                return {
                    'nodes': [(record['label'], record['count']) for record in summary],
                    'relationships': [(record['type'], record['count']) for record in relationships]
                }
        except Exception as e:
            st.error(f"Error getting graph summary: {str(e)}")
            return None

class RAGSystem:
    def __init__(self, system_type="traditional"):
        self.system_type = system_type
        self.llm = OpenAI(temperature=0)
        self.embeddings = OpenAIEmbeddings()
        self.vector_store = None
        
        if system_type == "kg":
            self.kg_manager = KnowledgeGraphManager(
                uri=os.getenv("NEO4J_URI"),
                user=os.getenv("NEO4J_USERNAME"),
                password=os.getenv("NEO4J_PASSWORD")
            )
            self.kg_manager.clear_database()

    def process_document(self, chunks):
        self.vector_store = FAISS.from_texts(chunks, self.embeddings)
        
        if self.system_type == "kg":
            self.build_knowledge_graph(chunks)

    def build_knowledge_graph(self, chunks):
        extract_prompt = PromptTemplate(
            input_variables=["text"],
            template="""
            Extract entities and their relationships from the following text.
            Return them in JSON format with 'entities' and 'relationships' keys.
            Each entity should have 'name' and 'type'.
            Each relationship should have 'source', 'target', and 'type'.
            Text: {text}
            """
        )
        
        chain = LLMChain(llm=self.llm, prompt=extract_prompt)
        
        for chunk in chunks:
            try:
                response = chain.invoke({"text": chunk})
                extracted_data = json.loads(response['text'])
                
                with self.kg_manager.driver.session() as session:
                    for entity in extracted_data['entities']:
                        session.run(
                            f"MERGE (n:{entity['type']} {{name: $name}})",
                            name=entity['name']
                        )
                    
                    for rel in extracted_data['relationships']:
                        session.run("""
                            MATCH (a), (b)
                            WHERE a.name = $source AND b.name = $target
                            MERGE (a)-[r:%s]->(b)
                        """ % rel['type'], source=rel['source'], target=rel['target'])
                        
            except Exception as e:
                st.warning(f"Error processing chunk: {str(e)}")

    def query(self, question):
        if not self.vector_store:
            return ("Please process a document first.", 0.0)

        start_time = time.time()
        
        try:
            if self.system_type == "traditional":
                qa_chain = RetrievalQA.from_chain_type(
                    llm=self.llm,
                    chain_type="stuff",
                    retriever=self.vector_store.as_retriever()
                )
                response = qa_chain.invoke({"query": question})
                answer = response['result'] if isinstance(response, dict) and 'result' in response else str(response)
            else:
                answer = self._kg_enhanced_query(question)

            query_time = time.time() - start_time
            return (answer, query_time)
            
        except Exception as e:
            return (f"Error processing query: {str(e)}", 0.0)

    def _kg_enhanced_query(self, question):
        relevant_docs = self.vector_store.similarity_search(question, k=3)
        vector_context = "\n".join([doc.page_content for doc in relevant_docs])
        
        kg_context = self._get_relevant_kg_context(question)
        
        enhanced_prompt = PromptTemplate(
            input_variables=["question", "vector_context", "kg_context"],
            template="""
            Answer the question using both the context and knowledge graph relationships.
            
            Context:
            {vector_context}
            
            Knowledge Graph Context:
            {kg_context}
            
            Question: {question}
            Answer:
            """
        )
        
        chain = LLMChain(llm=self.llm, prompt=enhanced_prompt)
        response = chain.invoke({
            "question": question,
            "vector_context": vector_context,
            "kg_context": kg_context
        })
        
        return response['text']

    def _get_relevant_kg_context(self, question):
        extract_terms_prompt = PromptTemplate(
            input_variables=["question"],
            template="Extract key terms from this question: {question}\nReturn only the terms, separated by commas."
        )
        terms_chain = LLMChain(llm=self.llm, prompt=extract_terms_prompt)
        terms_response = terms_chain.invoke({"question": question})
        key_terms = [term.strip() for term in terms_response['text'].split(',')]

        with self.kg_manager.driver.session() as session:
            relationships = []
            for term in key_terms:
                result = session.run("""
                    MATCH (n)-[r]-(m)
                    WHERE toLower(n.name) CONTAINS toLower($term)
                       OR toLower(m.name) CONTAINS toLower($term)
                    RETURN n.name as source, type(r) as relationship, m.name as target
                    LIMIT 5
                """, term=term)
                
                for record in result:
                    relationships.append(
                        f"{record['source']} {record['relationship']} {record['target']}"
                    )
            
            return "\n".join(relationships) if relationships else "No relevant relationships found."

def initialize_session_state():
    if 'is_processed' not in st.session_state:
        st.session_state.is_processed = False
    if 'processed_files' not in st.session_state:
        st.session_state.processed_files = []
    if 'current_status' not in st.session_state:
        st.session_state.current_status = "Waiting for documents"
    if 'rag_systems' not in st.session_state:
        st.session_state.rag_systems = None

def display_main_status():
    status_color = {
        "Waiting for documents": "#FFA500",
        "Processing": "#3498db",
        "Completed": "#4CAF50",
        "Error": "#ff4444"
    }
    color = status_color.get(st.session_state.current_status, "#FFA500")
    
    st.markdown(f"""
        <div class="main-status">
            <h3 style="color: {color}">Current Status: {st.session_state.current_status}</h3>
            <p style="color: white;">{'✅ Ready to answer questions' if st.session_state.is_processed else '⏳ Please upload and process documents'}</p>
        </div>
    """, unsafe_allow_html=True)



def initialize_session_state():
    """Initialize all session state variables"""
    if 'processed_files' not in st.session_state:
        st.session_state.processed_files = []
    if 'is_processed' not in st.session_state:
        st.session_state.is_processed = False
    if 'rag_systems' not in st.session_state:
        st.session_state.rag_systems = None
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = []
    if 'current_status' not in st.session_state:
        st.session_state.current_status = "Waiting for documents"

def display_chat_history():
    """Display chat history with proper error handling"""
    if not st.session_state.chat_history:
        st.info("No chat history yet. Start asking questions!")
        return

    for i, (question, trad_response, kg_response) in enumerate(st.session_state.chat_history):
        st.markdown("---")
        st.markdown(f"### Question {i+1}")
        st.info(f"Q: {question}")
        
        col1, col2 = st.columns(2)
        with col1:
            st.markdown("#### Traditional RAG Response")
            st.success(trad_response[0])
            st.metric("Response Time", f"{trad_response[1]:.2f}s")
        
        with col2:
            st.markdown("#### KG-Enhanced RAG Response")
            st.success(kg_response[0])
            st.metric("Response Time", f"{kg_response[1]:.2f}s")

def main():
    # Initialize session state at the start
    initialize_session_state()
    
    st.title("🤖 RAG vs Knowledge Graph-Enhanced RAG")
    
    # Display current status
    st.markdown(f"""
        <div class="main-status">
            <h3>Current Status: {st.session_state.current_status}</h3>
            <p>{'✅ Ready to answer questions' if st.session_state.is_processed else '⏳ Please upload and process documents'}</p>
        </div>
    """, unsafe_allow_html=True)

    # File upload section
    uploaded_files = st.file_uploader(
        "Upload your documents",
        type=['pdf', 'docx', 'txt', 'pptx'],
        accept_multiple_files=True
    )

    if uploaded_files:
        col1, col2 = st.columns([4, 1])
        with col2:
            process_button = st.button("🚀 Process Documents")

        if process_button:
            try:
                # Initialize processors
                doc_processor = DocumentProcessor()
                traditional_rag = RAGSystem(system_type="traditional")
                kg_rag = RAGSystem(system_type="kg")

                progress_bar = st.progress(0)
                status_text = st.empty()

                # Process documents
                all_chunks = []
                for idx, file in enumerate(uploaded_files):
                    status_text.text(f"Processing {file.name}...")
                    chunks = doc_processor.process_file(file)
                    all_chunks.extend(chunks)
                    
                    progress = (idx + 1) / len(uploaded_files)
                    progress_bar.progress(progress)

                # Process all chunks together
                if all_chunks:
                    traditional_rag.process_document(all_chunks)
                    kg_rag.process_document(all_chunks)
                    
                    st.session_state.is_processed = True
                    st.session_state.rag_systems = (traditional_rag, kg_rag)
                    st.session_state.processed_files = [f.name for f in uploaded_files]
                    st.session_state.current_status = "Completed"
                    
                    status_text.text("✅ Documents processed successfully!")
                    st.success("Ready to answer questions!")
                else:
                    st.error("No content could be extracted from the documents.")

            except Exception as e:
                st.error(f"Error during processing: {str(e)}")
                st.session_state.current_status = "Error"

    # Query section
    if st.session_state.is_processed and st.session_state.rag_systems:
        st.markdown("### 🔍 Ask Questions")
        query = st.text_input("Enter your question:")

        if query and st.button("Compare RAG vs KG-RAG"):
            traditional_rag, kg_rag = st.session_state.rag_systems
            
            with st.spinner("Processing query..."):
                # Get responses
                trad_response = traditional_rag.query(query)
                kg_response = kg_rag.query(query)
                
                # Add to chat history
                st.session_state.chat_history.append((query, trad_response, kg_response))
                
                # Display current response
                col1, col2 = st.columns(2)
                with col1:
                    st.markdown("### Traditional RAG")
                    st.write(trad_response[0])
                    st.metric("Response Time", f"{trad_response[1]:.2f}s")
                
                with col2:
                    st.markdown("### KG-Enhanced RAG")
                    st.write(kg_response[0])
                    st.metric("Response Time", f"{kg_response[1]:.2f}s")
                
                # Display Knowledge Graph
                # st.markdown("### 🕸️ Knowledge Graph Visualization")
                # graph_html = kg_rag.kg_manager.create_graph_visualization()
                # if graph_html:
                #     st.components.v1.html(graph_html, height=600)
                    
                    # Display graph summary
                    summary = kg_rag.kg_manager.get_graph_summary()
                    if summary:
                        col1, col2 = st.columns(2)
                        with col1:
                            st.markdown("#### Entity Types")
                            for label, count in summary['nodes']:
                                st.write(f"- {label[0]}: {count}")
                        with col2:
                            st.markdown("#### Relationship Types")
                            for rel_type, count in summary['relationships']:
                                st.write(f"- {rel_type}: {count}")

        # Display chat history
        if st.session_state.chat_history:
            st.markdown("### 💬 Chat History")
            display_chat_history()

if __name__ == "__main__":
    main() 


